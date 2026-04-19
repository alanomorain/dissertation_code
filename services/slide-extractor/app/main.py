from __future__ import annotations

import io
import re
from typing import Any

import fitz
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import JSONResponse
from pptx import Presentation

MAX_UPLOAD_BYTES = 10 * 1024 * 1024
ALLOWED_MIME = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

app = FastAPI(title="slide-extractor", version="1.0.0")


class ExtractorError(Exception):
    def __init__(self, code: str, message: str, status: int = 400):
        self.code = code
        self.message = message
        self.status = status
        super().__init__(message)


def normalize_text(value: str) -> str:
    lines = [re.sub(r"\s+", " ", line).strip() for line in value.splitlines()]

    normalized: list[str] = []
    previous_blank = False
    for line in lines:
        if not line:
            if not previous_blank:
                normalized.append("")
            previous_blank = True
            continue

        normalized.append(line)
        previous_blank = False

    return "\n".join(normalized).strip()


def extract_pdf_chunks(content: bytes) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    with fitz.open(stream=content, filetype="pdf") as doc:
        for page_number, page in enumerate(doc, start=1):
            text = normalize_text(page.get_text("text") or "")
            chunks.append(
                {
                    "index": page_number,
                    "label": f"Page {page_number}",
                    "text": text,
                }
            )
    return chunks


def shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""

    frame = shape.text_frame
    if frame is None:
        return ""

    return frame.text or ""


def extract_pptx_chunks(content: bytes) -> list[dict[str, Any]]:
    presentation = Presentation(io.BytesIO(content))
    chunks: list[dict[str, Any]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        pieces: list[str] = []

        for shape in slide.shapes:
            text = shape_text(shape).strip()
            if text:
                pieces.append(text)

        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()

        body = "\n".join(pieces)
        if notes_text:
            body = f"{body}\n\nSpeaker notes:\n{notes_text}" if body else f"Speaker notes:\n{notes_text}"

        text = normalize_text(body)
        chunks.append(
            {
                "index": slide_number,
                "label": f"Slide {slide_number}",
                "text": text,
            }
        )

    return chunks


def detect_kind(file: UploadFile) -> str:
    content_type = (file.content_type or "").strip().lower()
    filename = (file.filename or "").strip().lower()

    if content_type in ALLOWED_MIME:
        return ALLOWED_MIME[content_type]

    if filename.endswith(".pdf"):
        return "pdf"
    if filename.endswith(".pptx"):
        return "pptx"
    if filename.endswith(".ppt"):
        raise ExtractorError(
            code="UNSUPPORTED_TYPE",
            message="Legacy .ppt files are not supported. Please upload PDF or PPTX.",
            status=415,
        )

    raise ExtractorError(
        code="UNSUPPORTED_TYPE",
        message="Unsupported file type. Please upload PDF or PPTX.",
        status=415,
    )


def build_warnings(chunks: list[dict[str, Any]], full_text: str) -> list[str]:
    warnings: list[str] = []

    empty_count = sum(1 for chunk in chunks if not chunk.get("text"))
    if empty_count > 0:
        warnings.append(f"{empty_count} slide/page entries contained no extractable text.")

    if len(full_text) < 200:
        warnings.append("Very little text was extracted; topic quality may be lower.")

    return warnings


@app.exception_handler(ExtractorError)
async def extractor_error_handler(_request, exc: ExtractorError):
    return JSONResponse(
        status_code=exc.status,
        content={"error": {"code": exc.code, "message": exc.message}},
    )


@app.get("/health")
async def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/extract")
async def extract(
    file: UploadFile = File(...),
    moduleCode: str | None = Form(default=None),
):
    try:
        kind = detect_kind(file)
        content = await file.read()

        if not content:
            raise ExtractorError("EXTRACTION_FAILED", "The uploaded file is empty.", 400)

        if len(content) > MAX_UPLOAD_BYTES:
            raise ExtractorError("EXTRACTION_FAILED", "File is too large. Maximum allowed size is 10MB.", 413)

        if kind == "pdf":
            chunks = extract_pdf_chunks(content)
        else:
            chunks = extract_pptx_chunks(content)

        full_text = normalize_text("\n\n".join(chunk.get("text", "") for chunk in chunks if chunk.get("text")))

        if not full_text:
            raise ExtractorError(
                "EMPTY_EXTRACTION",
                "No extractable text found in this file. Please upload a text-based PDF/PPTX.",
                422,
            )

        warnings = build_warnings(chunks, full_text)

        return {
            "full_text": full_text,
            "slide_chunks": chunks,
            "stats": {
                "entry_count": len(chunks),
                "character_count": len(full_text),
            },
            "warnings": warnings,
            "meta": {
                "module_code": moduleCode or "",
                "filename": file.filename or "",
                "kind": kind,
            },
        }
    except ExtractorError:
        raise
    except Exception:
        raise ExtractorError(
            "EXTRACTION_FAILED",
            "Unable to extract text from this file.",
            500,
        )
